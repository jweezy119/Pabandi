#!/usr/bin/env python3
"""
Build Pabandi's NIC Karachi Cohort 16 Pitch Deck from the official template.
Preserves template layout/design, replaces placeholder text with Pabandi content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import re
import sys
import os

TEMPLATE = "/home/peesee/Pabandi/NIC-KHI-Pitch-Deck-Template-2026.pptx"
OUTPUT = "/home/peesee/Pabandi/Pabandi-NIC-Karachi-Cohort16.pptx"

# ── Slide Content Definitions ─────────────────────────────────────────────

SLIDE_CONTENT = {
    # Slide 1: Title & Logo
    1: {
        "Title & Logo of": "Pabandi",
        "Company": "",
        "One liner description of what you do": "The AI-Powered Trust Protocol & Escrow Layer for the Informal Economy",
    },
    # Slide 2: Vision
    2: {
        "Why does your company exists?": "Why Pabandi Exists",
        "How will you transform the industry or vertical you're operating in?": "Pakistan's informal economy — salons, live sellers, freelancers, property managers — runs on trust but has zero infrastructure to measure, verify, or reward it. Every day, businesses lose revenue to no-shows, fake payment screenshots, and COD fraud. Customers lose money to unverified sellers. There is no portable trust layer.",
        "What positive impact will your product(s) and solution(s)": "Pabandi transforms this by making trust visible, portable, and fair. We are building",
        "create for the world at large?": "the reliability infrastructure that gives honest people a measurable advantage — reducing fraud, eliminating no-shows, and enabling Pakistan's 20M+ freelancers and 500K+ service businesses to carry their reputation across every platform they touch.",
        "Vision": "Vision",
    },
    # Slide 3: Traction
    3: {
        "If you are seasoned professionals, then what your notable achievements to date?": "Founder: 8+ years IT infrastructure & automation (M365, PowerShell, REST APIs, cloud). Self-taught blockchain, ML, and full-stack development. Built entire Pabandi platform solo with AI coding agents.",
        "If you are a established company then how products have built/ launched so far?": "Full-stack production app: 76 client pages, 90+ backend services, 71 API routes, 5 Solidity smart contracts. Live on Firebase + Google Cloud Run.",
        "What are the milestones that your product have achieved in terms of:": "Milestones achieved:",
        "Revenue:": "Revenue model validated:",
        "to date": "$29/$79/$249 subscription tiers + API metering ($0.02-$0.05/call) + 2% escrow fee. Treasury orchestrator validated with $40K simulated flow.",
        "Users Acquired:": "Users/Integrations:",
        # the second "to date" after users
        "Users testimonials about the product:": "Alibaba CoCreate 2026 official entrant. Hub71 application drafted. Full Pakistan regulatory compliance architecture (SBP, PVARA, SECP). $PAB token live on Solana mainnet — non-mintable, non-freezable, LP burned.",
        "(state here)": "",
        "Any other achievements or accolades:": "Achievements:",
    },
    # Slide 4: Problem
    4: {
        "What specific problem are you solving for your customer?": "The Trust Crisis in Pakistan's Informal Economy",
        "Describe the pain or problem that you are alleviating or the pleasure you are providing to improve the quality of life of your customer/ client/ beneficiary.": "• 20-30% no-show rates destroy salon, restaurant, and clinic revenue nationwide\n• Fake payment screenshots on WhatsApp commerce cost sellers billions annually\n• COD fraud on Daraz/e-commerce — 16-30% rejection rates bleed merchant margins\n• No portable reputation — a reliable customer is treated identically to a ghost\n• Blanket deposits insult loyal customers and reduce first-time bookings by 40%\n• Fake reviews on OTAs and marketplaces make trust impossible to verify\n• Pakistan's $15B informal economy has ZERO trust infrastructure",
    },
    # Slide 5: Solution
    5: {
        "What is the minimum viable product (MVP) that embodies your UVP?": "Pabandi: Trust Made Portable, Deposits Made Fair",
        "What your product or service is and how it works in three simple steps.": "Step 1: SCORE — AI analyzes booking history, OSINT footprint, behavioral biometrics, and on-chain activity to generate a dynamic Trust Score (0-100) across 4 verticals (Commerce, Hospitality, Freelance, Appointment).\n\nStep 2: PROTECT — Smart escrow locks funds on booking. High-trust users pay zero deposit. New/risky users pay fair, proportional deposits. Never more than 50%. Always explained.\n\nStep 3: REWARD — Show up = earn $PAB tokens (Sharia-compliant Hibah rewards). Score improves. Deposit drops. Soulbound NFT badge minted on-chain. Trust becomes portable.",
        "Show screenshots/mockups of the product/ solution": "Live at pabandi.com — Production-deployed on Firebase + Google Cloud Run",
    },
    # Slide 6: Tech Stack & Roadmap
    6: {
        "How technology is enabling your business?": "Full-Stack AI + Web3 Trust Infrastructure",
        "Share brief details of your Tech Stack and Architecture": "Frontend: React 18 + TypeScript + TailwindCSS + Vite + Zustand\nBackend: Node.js + Express + Prisma ORM + PostgreSQL\nBlockchain: Multi-chain (Solana, ETH, BNB, BTC, Stellar) — 5 smart contracts\nAI/ML: 4-model ensemble (GBT, Temporal GNN, Wide&Deep NN, Meta-Learner) + OSINT engines (Threat Fusion, Adversarial Graph, Temporal Deception, Behavioral Biometrics)\nProtocol: Pabandi Trust Protocol (PTP) — open standard for portable trust attestations with ZK proofs\nIntegrations: Shopify SDK, WhatsApp Evolution API, Channex/Beds24/Cloudbeds PMS\nInfra: Firebase Hosting, Google Cloud Run, Docker\n\nRoadmap:\nQ3 2026: NIC incubation + first 100 Pakistani merchants onboarded\nQ4 2026: Daraz seller integration pilot + Shopify App Store launch\nQ1 2027: PVARA VASP licensing + SDK expansion (Python, React Native)\nQ2 2027: Cross-platform Passport (5 integrations) + Indonesia/Nigeria pilot",
        "(Do bring the Hardware, in case the product/solution is Hardware based)": "",
    },
    # Slide 7: Customer Segments
    7: {
        "Who are the specific groups of people experiencing this problem?": "Four Primary Customer Segments",
        "Explain your customer persona.": "Segment 1 — Service SMEs (Salons, Clinics, Restaurants): Pakistani business owners losing 20-30% revenue to no-shows. Need: risk-based deposits, reliable booking system. Willingness to pay: $29-$79/mo.\n\nSegment 2 — Live Sellers & Social Commerce: Instagram/TikTok/WhatsApp sellers losing to fake orders and payment fraud. Need: escrow checkout links, buyer verification. 500K+ active in Pakistan alone.\n\nSegment 3 — E-Commerce Platforms (Daraz, Shopify stores): Losing to COD fraud, brushing scams, fake reviews. Need: API-based reliability scoring. Willingness to pay: $0.02-$0.05/API call.\n\nSegment 4 — Freelancers & Gig Workers: 20M+ Pakistani freelancers with no portable reputation. Need: cross-platform Trust Passport, escrow-backed contracts. Value: reduced client ghosting, faster payments.",
    },
    # Slide 8: Market Sizing
    8: {
        "Total Addressable Market": "Total Addressable Market",
        "TAM": "$48B — Global informal economy trust infrastructure. Hotel no-shows ($12-15B/yr), e-commerce fraud ($41B/yr), gig economy trust gap.",
        "Serviceable Attainable Market": "Serviceable Addressable Market",
        "SAM": "$2.4B — Pakistan + MENA informal commerce trust layer. 500K service businesses, 20M freelancers, $15B informal economy in Pakistan alone.",
        "Serviceable Obtainable Market": "Serviceable Obtainable Market",
        "SOM": "$12M — First 3 years. 5,000 merchant subscriptions ($79/mo avg) + 50M API calls/yr ($0.02/call) + 2% escrow fees on $30M transaction volume.",
    },
    # Slide 9: UVP
    9: {
        "What unique solution do you offer to the identified problem?": "The Only Sharia-Compliant Trust Protocol with Financial Guarantees",
        "Provide a Single, clear, compelling message that states \"why\" you are different and worth paying attention.": "Pabandi is the ONLY platform that combines:\n\n✦ AI-powered multi-axis trust scoring (not just reviews — OSINT + behavioral biometrics + on-chain identity)\n✦ Smart escrow that adapts to trust level (zero deposit for reliable users, fair deposits for unknowns)\n✦ An open trust protocol (PTP) with cryptographic attestations any platform can verify — even offline\n✦ Financial guarantees — Band A/B merchants are insured up to $500/transaction\n✦ Sharia-compliant rewards (Mudarabah profit-sharing, Hibah gifts — zero Riba)\n✦ Multi-vertical scoring — separate scores for Commerce, Hospitality, Freelance, Appointments\n✦ Pakistan-first architecture with global scalability\n\nWada pura karo. Inaam pao.\n(Keep your word. Get rewarded.)",
    },
    # Slide 10: Competitor Analysis
    10: {
        "Who are your 5 biggest competitors (nationally or internationally)?": "Competitive Landscape — No Direct Competitor Exists",
        "What are the 5 alternatives available in the market that your prospect customers/clients/users are using to solve their problem(s)?": "1. Binance P2P / Bybit — Manual C2C arbitration for crypto-fiat. Slow, no trust scoring, not Sharia-compliant. Pabandi: 0-second atomic settlement + AI trust.\n\n2. Daraz Seller Center — Basic seller ratings, no portable trust, rampant COD fraud and brushing scams. Pabandi: OSINT-powered fraud detection API.\n\n3. Depositly (Canada) — Deposit management for formal businesses only. No AI scoring, no crypto, no Pakistan presence. Pabandi: AI + Web3 + informal economy focus.\n\n4. Trustpilot / Google Reviews — Fakeable, no proof-of-transaction, no financial guarantees. Pabandi: on-chain verified reviews + Soulbound identity.\n\n5. Manual WhatsApp + Cash — How 90% of Pakistan's informal economy operates today. Zero protection, zero portability, zero rewards. Pabandi replaces this entirely.",
    },
    # Slide 11: Channels
    11: {
        "How will you reach your target customer segments?": "Multi-Channel Distribution Strategy",
        "Identify the marketing channels from which you can capture customers/ clients / users.": "1. WhatsApp-First Outreach — AI-powered WhatsApp bot with Evolution API gateway. Merchants receive booking links directly in their existing WhatsApp business flows. Zero app download required.\n\n2. Shopify App Store — Zero-config widget SDK. Merchants paste one JS snippet to add Trust Badge + Escrow Checkout to their stores.\n\n3. Content & Social — Roman Urdu content engine (Twitter/X, Instagram, TikTok). 90-day content calendar already built. Viral social proof receipts shared by buyers.\n\n4. B2B API Partnerships — Direct integration with Daraz, TikTok Shop, and local POS systems. Developer portal with Swagger docs live.\n\n5. Referral Flywheel — 500 $PAB per referral. Every successful checkout generates a shareable receipt card for WhatsApp — organic viral loop.\n\n6. NIC Karachi Network — Leverage NIC's 1000+ alumni network for first 100 merchant pilots.",
    },
    # Slide 12: Revenue Streams
    12: {
        "How will you generate revenue from your solution?": "Five Revenue Streams — Never from Selling Data",
        "Revenue Model/List your sources of revenue:": "Revenue Model:",
        "For example:": "",
        "Subscription based, transaction based, project based, services based.": "1. SaaS Subscriptions (Recurring): $29/mo Starter | $79/mo Growth | $249/mo Enterprise. Dashboard, analytics, risk deposits, CRM, WhatsApp automation.\n\n2. API Metering (Usage-Based): $0.02-$0.05 per trust assessment call. Developers and platforms pay per verification. 10x margin on ~$0.0001 server cost.\n\n3. Transaction Fees: 2% platform fee on all escrow deposits. Scales linearly with volume.\n\n4. White-Label Licensing: $5K-$25K/year for custom-branded trust infrastructure. Hotel chains, franchise networks, enterprise clients.\n\n5. Halal Staking Commission: 0.5% on Mudarabah profit-share pool. Revenue grows with staking volume.\n\nZero revenue from: selling user data, behavioral advertising, or surveillance capitalism.",
    },
    # Slide 13: Cost Structure
    13: {
        "What are your most important costs associated with developing and delivering your solution?": "Lean Cost Structure — Software-First, Pakistan-Optimized",
        "Capital Expenditure - CAPEX:": "CAPEX:",
        "Operational Expenditure - OPEX:": "OPEX:",
    },
    # Slide 14: Financials
    14: {
        "Please provide current or forecasted financials:": "Forecasted Financials (Year 1-3)",
        "Balance Sheet /": "",
        "P&L Statement/": "",
        "Cash flow statement": "",
    },
    # Slide 15: Exit / Ask
    15: {
        "What's your ask and against what equity? (How much funding are you looking for to scale your startup?)": "The Ask: NIC Incubation + Pre-Seed Readiness",
        "What are your plans for IPO and/or exit?": "Immediate Ask: NIC Karachi Cohort 16 incubation — mentorship, workspace, regulatory guidance, and network access to onboard first 100 Pakistani merchants.\n\nPre-Seed Target: $150K for 12-month runway. Use: PVARA VASP licensing ($30K), first BD hire ($40K), merchant acquisition ($40K), infrastructure ($20K), legal/compliance ($20K).\n\nExit Strategy:\n• Year 3-5: Strategic acquisition target for Daraz/Alibaba (trust infrastructure for their ecosystem), Careem (trust layer for MENA services), or JazzCash/Easypaisa (trust-backed payments).\n• Year 5-7: Series A/B path toward regional IPO (PSX or ADX) as Pakistan's first trust infrastructure company.\n• Alternative: Open-source the PTP protocol, monetize enterprise licensing globally (the 'Red Hat model' for trust).",
    },
    # Slide 16: Key Metrics — this one has a table structure, handle carefully
    16: {
        "Customer Acquisition Cost - CAC:": "CAC: $0.50 (organic, content-driven) → $2.00 (paid channels at scale)",
        "Quarterly Recurring Revenue (QRR):": "QRR: $0 → $2.1K → $7.5K → $18.7K → $37.5K → $75K",
        "Life TIme Value of Customer LTV:": "LTV: $948 (avg $79/mo × 12mo retention) — LTV:CAC ratio = 474:1",
    },
    # Slide 17: Team
    17: {
        "The Visionary": "The Visionary",
        "(CEO)": "(Founder & CEO)",
        "The Hustler": "The Hustler",
        "(The person who manages": "(Operations &",
        "the operations)": "BD Lead)",
        "The Hacker": "The Hacker",
        "(the programmer or main": "(Full-Stack",
        "product developer)": "& AI/Web3)",
    },
    # Slide 18: Contact
    18: {
        "Website": "pabandi.com",
        "Email id": "team@pabandi.com",
        "Address": "Karachi, Pakistan (NIC Incubation Target)",
        "Cell phone -": "Contact:",
        "Landline": "[Phone Number]",
    },
}


def replace_text_in_shape(shape, replacements):
    """Replace text in a shape's text frame, preserving formatting."""
    if not shape.has_text_frame:
        return False
    
    changed = False
    for paragraph in shape.text_frame.paragraphs:
        full_para_text = ''.join(run.text for run in paragraph.runs)
        
        for old_text, new_text in replacements.items():
            if old_text in full_para_text:
                # If the paragraph has runs, try to replace in the first run
                # and clear subsequent runs that were part of the old text
                if len(paragraph.runs) == 1:
                    paragraph.runs[0].text = paragraph.runs[0].text.replace(old_text, new_text)
                    changed = True
                elif len(paragraph.runs) > 1:
                    # Rebuild: put all text in first run, clear rest
                    new_full = full_para_text.replace(old_text, new_text)
                    paragraph.runs[0].text = new_full
                    for run in paragraph.runs[1:]:
                        run.text = ""
                    changed = True
                elif len(paragraph.runs) == 0 and paragraph.text:
                    # Direct text without runs
                    pass
    
    return changed


def process_slide(slide, slide_num, content_map):
    """Process a single slide, replacing placeholder text."""
    if slide_num not in content_map:
        return
    
    replacements = content_map[slide_num]
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_in_shape(shape, replacements)
        
        # Handle grouped shapes
        if shape.shape_type == 6:  # Group shape
            try:
                for child_shape in shape.shapes:
                    if child_shape.has_text_frame:
                        replace_text_in_shape(child_shape, replacements)
            except Exception:
                pass
        
        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        full_text = ''.join(run.text for run in paragraph.runs)
                        for old_text, new_text in replacements.items():
                            if old_text in full_text:
                                new_full = full_text.replace(old_text, new_text)
                                if paragraph.runs:
                                    paragraph.runs[0].text = new_full
                                    for run in paragraph.runs[1:]:
                                        run.text = ""


def main():
    print(f"Loading template: {TEMPLATE}")
    prs = Presentation(TEMPLATE)
    
    print(f"Template has {len(prs.slides)} slides")
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {i}...")
        process_slide(slide, i, SLIDE_CONTENT)
    
    # Save
    prs.save(OUTPUT)
    print(f"\n✅ Pitch deck saved to: {OUTPUT}")
    print(f"   File size: {os.path.getsize(OUTPUT) / 1024:.1f} KB")


if __name__ == "__main__":
    main()
